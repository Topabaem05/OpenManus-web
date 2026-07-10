from __future__ import annotations

import os
from pathlib import Path
from typing import Any

from app.config import config
from app.llm import LLM
from app.tool.base import BaseTool, ToolResult


class DocumentGeneration(BaseTool):
    """Generate documents (docx, md) and slide decks (pptx) from a structured outline."""

    name: str = "document_generation"
    description: str = (
        "Generate a document or slide deck from a structured outline. "
        "Supported formats: markdown (.md), Word (.docx), PowerPoint (.pptx). "
        "Input: title, format, sections (list of {heading, bullets, level})."
    )
    parameters: dict = {
        "type": "object",
        "properties": {
            "title": {"type": "string", "description": "Document or deck title"},
            "format": {
                "type": "string",
                "enum": ["markdown", "docx", "pptx"],
                "default": "markdown",
            },
            "sections": {
                "type": "array",
                "items": {
                    "type": "object",
                    "properties": {
                        "heading": {"type": "string"},
                        "bullets": {"type": "array", "items": {"type": "string"}},
                        "level": {"type": "integer", "default": 1},
                        "slide_layout": {
                            "type": "string",
                            "enum": ["title", "title_and_content"],
                            "default": "title_and_content",
                        },
                    },
                    "required": ["heading"],
                },
            },
            "output_path": {
                "type": "string",
                "description": "Relative path under workspace/; defaults to title + format",
            },
        },
        "required": ["title", "sections"],
    }

    def __init__(self, **data: Any):
        super().__init__(**data)
        self._llm_cache: LLM | None = None

    def _llm(self) -> LLM:
        if self._llm_cache is None:
            self._llm_cache = LLM(config_name="document_generation")
        return self._llm_cache

    async def execute(self, **kwargs: Any) -> ToolResult:
        title = kwargs.get("title", "Untitled")
        fmt = kwargs.get("format", "markdown")
        sections = kwargs.get("sections", []) or []
        output_path = kwargs.get("output_path", "")
        if not output_path:
            safe_title = "".join(c if c.isalnum() or c in "-_ " else "_" for c in title)
            output_path = f"{safe_title.replace(' ', '_')}.{fmt if fmt != 'markdown' else 'md'}"

        workspace = Path(config.workspace_root)
        out = workspace / output_path
        out.parent.mkdir(parents=True, exist_ok=True)

        if fmt == "pptx":
            return await self._make_pptx(title, sections, out)
        if fmt == "docx":
            return await self._make_docx(title, sections, out)
        return await self._make_markdown(title, sections, out)

    async def _make_markdown(self, title: str, sections: list, out: Path) -> ToolResult:
        lines = [f"# {title}", ""]
        for section in sections:
            heading = section.get("heading", "")
            level = max(1, min(6, section.get("level", 1)))
            lines.append(f"{'#' * level} {heading}")
            for bullet in section.get("bullets", []) or []:
                lines.append(f"- {bullet}")
            lines.append("")
        out.write_text("\n".join(lines), encoding="utf-8")
        return ToolResult(output=f"Markdown saved to {out}")

    async def _make_docx(self, title: str, sections: list, out: Path) -> ToolResult:
        try:
            from docx import Document
            from docx.enum.text import WD_ALIGN_PARAGRAPH
        except Exception as exc:
            return ToolResult(error=f"python-docx not installed: {exc}")
        doc = Document()
        title_para = doc.add_heading(title, level=0)
        title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        for section in sections:
            heading = section.get("heading", "")
            level = section.get("level", 1)
            doc.add_heading(heading, level=max(1, min(6, level)))
            for bullet in section.get("bullets", []) or []:
                doc.add_paragraph(bullet, style="List Bullet")
            doc.add_paragraph()
        tmp = out.with_suffix(".tmp.docx")
        doc.save(str(tmp))
        tmp.rename(out)
        return ToolResult(output=f"Word document saved to {out}")

    async def _make_pptx(self, title: str, sections: list, out: Path) -> ToolResult:
        try:
            from pptx import Presentation
        except Exception as exc:
            return ToolResult(error=f"python-pptx not installed: {exc}")
        prs = Presentation()
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        if slide.placeholders:
            slide.placeholders[1].text = "Generated by OpenManus"
        for section in sections:
            heading = section.get("heading", "")
            bullets = section.get("bullets", []) or []
            layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = heading
            body = slide.shapes.placeholders[1]
            tf = body.text_frame
            tf.text = bullets[0] if bullets else ""
            for bullet in bullets[1:]:
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 0
        tmp = out.with_suffix(".tmp.pptx")
        prs.save(str(tmp))
        tmp.rename(out)
        return ToolResult(output=f"PowerPoint deck saved to {out}")
